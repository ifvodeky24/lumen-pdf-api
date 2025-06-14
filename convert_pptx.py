# convert_pptx.py
import sys
import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

pdf_file = sys.argv[1]
pptx_file = sys.argv[2]

# Konversi halaman PDF ke gambar
images = convert_from_path(pdf_file, dpi=200)

# Buat PPT
prs = Presentation()

# Set ukuran slide ke ukuran gambar
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

for image in images:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Simpan sementara sebagai JPG
    temp_img = 'temp_slide.jpg'
    image.save(temp_img, 'JPEG')
    
    # Masukkan ke slide
    slide.shapes.add_picture(temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    os.remove(temp_img)

prs.save(pptx_file)

print("Conversion successful")
